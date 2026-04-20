"""
html_to_pptx.py
───────────────
Convert HTML slides to a PowerPoint (.pptx) file.

Installation (uv)
    uv add playwright python-pptx chromium
    uv run playwright install

Typical usage (from a chat/LLM app that collects confirmed HTML per slide):

    from html_to_pptx import HtmlToPptxConverter

    converter = HtmlToPptxConverter()
    converter.add_slide_from_string(html1, title="Intro", notes="Say hello")
    converter.add_slide_from_string(html2, title="Architecture")
    converter.add_slide_from_file("slide3.html", title="Summary")
    converter.save("output.pptx")

Overflow safety:
    Each slide is checked for overflowing elements before the screenshot is
    taken. The default behaviour is to warn. Use on_overflow="raise" to treat
    any overflow as a hard error, or on_overflow="ignore" to skip checks.

    You can also call check_slide() / check_all_slides() explicitly to get a
    structured report before committing to save().

Pipeline per slide:
    HTML (string or file)
        → Playwright headless Chromium
            → overflow check (JS DOM inspection)
            → screenshot → PNG
        → python-pptx (full-bleed image on blank slide)
        → optional: hidden title placeholder + speaker notes
"""

from __future__ import annotations

import json
import os
import tempfile
import warnings
from dataclasses import dataclass, field
from pathlib import Path
from typing import Literal, Optional

from playwright.sync_api import sync_playwright, Page
from pptx import Presentation
from pptx.util import Inches


# ── Slide dimensions ──────────────────────────────────────────────────────────
# 16:9 widescreen: 10" × 5.625"  (standard modern PowerPoint)
SLIDE_WIDTH_INCH  = 10.0
SLIDE_HEIGHT_INCH = 5.625

# Viewport for the headless browser (pixels).
# 1280 × 720 maps cleanly to 10" × 5.625" @ 128 DPI.
VIEWPORT_W = 1280
VIEWPORT_H = 720

# Tolerance in px — elements 1-2 px outside the boundary are likely sub-pixel
# rounding artefacts, not real content overflow.
OVERFLOW_TOLERANCE_PX = 2

OnOverflow = Literal["warn", "raise", "ignore"]


# ── Data classes ──────────────────────────────────────────────────────────────

@dataclass
class OverflowViolation:
    """
    Describes a single DOM element that bleeds outside the slide canvas.

    Attributes
    ----------
    slide_index : int
        0-based index of the offending slide.
    tag : str
        HTML tag name of the element (e.g. "div", "p", "img").
    text_snippet : str
        First 80 characters of the element's text content (empty for images etc.).
    rect : dict
        Raw bounding rect: {"top", "right", "bottom", "left", "width", "height"}.
    overflow_edges : list[str]
        Which edges are violated, e.g. ["bottom", "right"].
    overflow_px : dict[str, float]
        How many pixels each edge overflows by, e.g. {"bottom": 34.5}.
    """
    slide_index: int
    tag: str
    text_snippet: str
    rect: dict
    overflow_edges: list[str]
    overflow_px: dict[str, float]

    def __str__(self) -> str:
        edges = ", ".join(
            f"{e} by {self.overflow_px[e]:.1f}px" for e in self.overflow_edges
        )
        snippet = (' \u201c' + self.text_snippet + '\u201d') if self.text_snippet else ""
        return (
            f"Slide {self.slide_index + 1}: <{self.tag}>{snippet} "
            f"overflows [{edges}]"
        )


@dataclass
class SlideCheckResult:
    """
    Result of an overflow check for one slide.

    Attributes
    ----------
    slide_index : int
    has_overflow : bool
    violations : list[OverflowViolation]
        Empty when has_overflow is False.
    """
    slide_index: int
    has_overflow: bool
    violations: list[OverflowViolation] = field(default_factory=list)

    def summary(self) -> str:
        if not self.has_overflow:
            return f"Slide {self.slide_index + 1}: ✓ No overflow detected."
        lines = [f"Slide {self.slide_index + 1}: ✗ {len(self.violations)} overflow violation(s):"]
        for v in self.violations:
            lines.append(f"  • {v}")
        return "\n".join(lines)


@dataclass
class SlideData:
    """Internal representation of one pending slide."""
    html: str
    title: Optional[str] = None
    notes: Optional[str] = None


# ── JS injected into every page to detect overflow ───────────────────────────
# Returns a JSON array of violating elements, or an empty array.
_OVERFLOW_CHECK_JS = """
([canvasW, canvasH, tolerancePx]) => {
    const violations = [];
    const all = document.querySelectorAll('*');

    for (const el of all) {
        const r = el.getBoundingClientRect();
        if (r.width === 0 && r.height === 0) continue;

        const edges = {};
        if (r.top    < -tolerancePx)          edges.top    = -r.top;
        if (r.left   < -tolerancePx)          edges.left   = -r.left;
        if (r.bottom > canvasH + tolerancePx) edges.bottom = r.bottom - canvasH;
        if (r.right  > canvasW + tolerancePx) edges.right  = r.right  - canvasW;

        if (Object.keys(edges).length === 0) continue;

        violations.push({
            tag:          el.tagName.toLowerCase(),
            text_snippet: (el.textContent || '').trim().slice(0, 80),
            rect: {
                top: r.top, right: r.right,
                bottom: r.bottom, left: r.left,
                width: r.width, height: r.height,
            },
            overflow_edges: Object.keys(edges),
            overflow_px:    edges,
        });
    }

    // De-duplicate: keep only the outermost element when a parent and its
    // children all overflow the same edges (avoids noisy child-level reports).
    const deduped = violations.filter((v, _, arr) => {
        return !arr.some(other =>
            other !== v &&
            other.rect.top    <= v.rect.top    &&
            other.rect.left   <= v.rect.left   &&
            other.rect.bottom >= v.rect.bottom &&
            other.rect.right  >= v.rect.right  &&
            JSON.stringify(other.overflow_edges) === JSON.stringify(v.overflow_edges)
        );
    });

    return JSON.stringify(deduped);
}
"""


class HtmlToPptxConverter:
    """
    Collect HTML slides one by one, then export them all as a .pptx file.

    Parameters
    ----------
    on_overflow : "warn" | "raise" | "ignore"
        What to do when a slide has content overflowing the canvas.
        - "warn"   (default) — print a warning and continue.
        - "raise"  — raise OverflowError and abort save().
        - "ignore" — skip overflow checks entirely.

    Methods
    -------
    add_slide_from_string(html, title, notes)
    add_slide_from_file(path, title, notes)
    remove_slide(index)
    reorder_slides(new_order)
    check_slide(index) -> SlideCheckResult
        Run overflow check on a single queued slide (opens its own browser).
    check_all_slides() -> list[SlideCheckResult]
        Run overflow check on all queued slides.
    save(output_path) -> str
        Render all slides, run overflow checks, write .pptx.
    """

    def __init__(
        self,
        width_inch: float = SLIDE_WIDTH_INCH,
        height_inch: float = SLIDE_HEIGHT_INCH,
        viewport_w: int = VIEWPORT_W,
        viewport_h: int = VIEWPORT_H,
        on_overflow: OnOverflow = "warn",
        overflow_tolerance_px: int = OVERFLOW_TOLERANCE_PX,
    ):
        self._slides: list[SlideData] = []
        self._width_inch  = width_inch
        self._height_inch = height_inch
        self._viewport_w  = viewport_w
        self._viewport_h  = viewport_h
        self._on_overflow = on_overflow
        self._tolerance   = overflow_tolerance_px

    # ── Public API ────────────────────────────────────────────────────────────

    def add_slide_from_string(
        self,
        html: str,
        title: Optional[str] = None,
        notes: Optional[str] = None,
    ) -> "HtmlToPptxConverter":
        """Add a slide from a raw HTML string. Returns self for chaining."""
        self._slides.append(SlideData(html=html, title=title, notes=notes))
        return self

    def add_slide_from_file(
        self,
        path: str | Path,
        title: Optional[str] = None,
        notes: Optional[str] = None,
    ) -> "HtmlToPptxConverter":
        """Add a slide from an HTML file. Returns self for chaining."""
        html = Path(path).read_text(encoding="utf-8")
        return self.add_slide_from_string(html, title=title, notes=notes)

    def remove_slide(self, index: int) -> "HtmlToPptxConverter":
        """Remove the slide at 0-based *index*. Returns self for chaining."""
        if not 0 <= index < len(self._slides):
            raise IndexError(
                f"Slide index {index} out of range (have {len(self._slides)} slides)."
            )
        self._slides.pop(index)
        return self

    def reorder_slides(self, new_order: list[int]) -> "HtmlToPptxConverter":
        """
        Reorder slides by index list.

        Example
        -------
        converter.reorder_slides([2, 0, 1])  # move slide 2 to the front
        """
        if sorted(new_order) != list(range(len(self._slides))):
            raise ValueError(
                f"new_order must be a permutation of 0..{len(self._slides)-1}, "
                f"got {new_order}"
            )
        self._slides = [self._slides[i] for i in new_order]
        return self

    def check_slide(self, index: int) -> SlideCheckResult:
        """
        Run an overflow check on a single queued slide.

        Opens its own browser session — useful for interactive validation
        during the LLM review loop before all slides are ready.

        Returns
        -------
        SlideCheckResult
        """
        if not 0 <= index < len(self._slides):
            raise IndexError(
                f"Slide index {index} out of range (have {len(self._slides)} slides)."
            )
        with sync_playwright() as pw:
            browser, context, page = self._make_browser(pw)
            try:
                result = self._check_page(page, self._slides[index], index)
            finally:
                page.close()
                context.close()
                browser.close()
        return result

    def check_all_slides(self) -> list[SlideCheckResult]:
        """
        Run overflow checks on every queued slide.

        Returns
        -------
        list[SlideCheckResult]
            One entry per slide, in order.
        """
        results: list[SlideCheckResult] = []
        with sync_playwright() as pw:
            browser, context, page = self._make_browser(pw)
            try:
                for i, slide in enumerate(self._slides):
                    results.append(self._check_page(page, slide, i))
            finally:
                page.close()
                context.close()
                browser.close()
        return results

    def save(self, output_path: str | Path) -> str:
        """
        Render all slides and write the .pptx file.

        Overflow behaviour is governed by the *on_overflow* constructor param:
        - "warn"   → warnings.warn() for each violating slide, then continue.
        - "raise"  → raises OverflowError listing all violations before writing.
        - "ignore" → skips overflow checks entirely.

        Returns
        -------
        str
            Resolved absolute path of the saved .pptx file.
        """
        if not self._slides:
            raise ValueError(
                "No slides added. Call add_slide_from_string() or "
                "add_slide_from_file() first."
            )

        output_path = Path(output_path).resolve()
        output_path.parent.mkdir(parents=True, exist_ok=True)

        with tempfile.TemporaryDirectory(prefix="html_to_pptx_") as tmpdir:
            png_paths = self._render_screenshots(tmpdir)
            self._build_pptx(png_paths, output_path)

        return str(output_path)

    # ── Slide count helpers ───────────────────────────────────────────────────

    def __len__(self) -> int:
        return len(self._slides)

    def __repr__(self) -> str:
        return f"HtmlToPptxConverter({len(self._slides)} slide(s) pending)"

    # ── Internal: browser factory ─────────────────────────────────────────────

    def _make_browser(self, pw):
        """Launch a headless Chromium browser and return (browser, context, page)."""
        browser = pw.chromium.launch(
            headless=True,
            args=[
                "--no-sandbox",
                "--disable-setuid-sandbox",
                "--disable-dev-shm-usage",
            ],
        )
        context = browser.new_context(
            viewport={"width": self._viewport_w, "height": self._viewport_h},
        )

        def _block_external(route, request):
            """Block external network requests so self-contained HTML is enforced."""
            if request.resource_type in ("image", "stylesheet", "font", "script"):
                url = request.url
                if url.startswith("http") and "data:" not in url:
                    route.abort()
                    return
            route.continue_()

        context.route("**/*", _block_external)
        page = context.new_page()
        return browser, context, page

    # ── Internal: overflow detection ──────────────────────────────────────────

    def _check_page(self, page: Page, slide: SlideData, index: int) -> SlideCheckResult:
        """
        Load *slide* HTML into *page* and run the JS overflow detector.
        Returns a SlideCheckResult — does NOT take a screenshot.
        """
        page.set_content(slide.html, wait_until="domcontentloaded")
        page.wait_for_timeout(300)

        raw_json: str = page.evaluate(
            _OVERFLOW_CHECK_JS,
            [self._viewport_w, self._viewport_h, self._tolerance],
        )
        raw: list[dict] = json.loads(raw_json)

        violations = [
            OverflowViolation(
                slide_index=index,
                tag=item["tag"],
                text_snippet=item["text_snippet"],
                rect=item["rect"],
                overflow_edges=item["overflow_edges"],
                overflow_px=item["overflow_px"],
            )
            for item in raw
        ]

        return SlideCheckResult(
            slide_index=index,
            has_overflow=bool(violations),
            violations=violations,
        )

    def _handle_overflow_results(self, results: list[SlideCheckResult]) -> None:
        """
        Apply the on_overflow policy to a batch of check results.
        Called after all checks are done, before screenshots are taken.
        """
        if self._on_overflow == "ignore":
            return

        failing = [r for r in results if r.has_overflow]
        if not failing:
            return

        if self._on_overflow == "warn":
            for r in failing:
                warnings.warn(r.summary(), stacklevel=4)

        elif self._on_overflow == "raise":
            lines = [
                f"{sum(len(r.violations) for r in failing)} overflow violation(s) "
                f"across {len(failing)} slide(s):\n"
            ]
            for r in failing:
                lines.append(r.summary())
            raise OverflowError("\n".join(lines))

    # ── Internal: Phase 1 – Playwright screenshots ───────────────────────────

    def _render_screenshots(self, tmpdir: str) -> list[str]:
        """
        Use headless Chromium to:
          1. Check each slide for overflow (respects on_overflow policy).
          2. Screenshot each slide to a PNG file.

        Returns a list of absolute PNG paths in slide order.
        """
        png_paths: list[str] = []
        check_results: list[SlideCheckResult] = []

        with sync_playwright() as pw:
            browser, context, page = self._make_browser(pw)

            try:
                for i, slide in enumerate(self._slides):
                    page.set_content(slide.html, wait_until="domcontentloaded")
                    page.wait_for_timeout(300)

                    # ── Overflow check (before screenshot) ──
                    if self._on_overflow != "ignore":
                        result = self._check_page(page, slide, i)
                        check_results.append(result)

                    # ── Screenshot ──
                    png_path = os.path.join(tmpdir, f"slide_{i:03d}.png")
                    page.screenshot(
                        path=png_path,
                        clip={
                            "x": 0, "y": 0,
                            "width": self._viewport_w,
                            "height": self._viewport_h,
                        },
                        type="png",
                    )
                    png_paths.append(png_path)

            finally:
                page.close()
                context.close()
                browser.close()

        # Apply policy after all slides are processed so "raise" gives a full
        # report of every problem slide, not just the first one.
        if check_results:
            self._handle_overflow_results(check_results)

        return png_paths

    # ── Internal: Phase 2 – python-pptx assembly ─────────────────────────────

    def _build_pptx(self, png_paths: list[str], output_path: Path) -> None:
        """Assemble the .pptx from PNG screenshots."""
        prs = Presentation()
        prs.slide_width  = Inches(self._width_inch)
        prs.slide_height = Inches(self._height_inch)

        blank_layout = self._get_blank_layout(prs)

        for slide_data, png_path in zip(self._slides, png_paths):
            slide = prs.slides.add_slide(blank_layout)

            slide.shapes.add_picture(
                png_path,
                left=Inches(0), top=Inches(0),
                width=Inches(self._width_inch),
                height=Inches(self._height_inch),
            )

            if slide_data.title:
                self._set_slide_title(slide, slide_data.title)

            if slide_data.notes:
                self._set_speaker_notes(slide, slide_data.notes)

        prs.save(str(output_path))

    @staticmethod
    def _get_blank_layout(prs: Presentation):
        layouts = prs.slide_layouts
        try:
            return layouts[6]   # index 6 = "Blank" in default theme
        except IndexError:
            return layouts[-1]

    @staticmethod
    def _set_slide_title(slide, title: str) -> None:
        """
        Write the title into the hidden title placeholder (off-slide).
        Readable by PowerPoint outline / search / accessibility tools.
        """
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.idx == 0:
                placeholder.text = title
                placeholder.left   = Inches(-10)
                placeholder.top    = Inches(-10)
                placeholder.width  = Inches(1)
                placeholder.height = Inches(0.1)
                return

    @staticmethod
    def _set_speaker_notes(slide, notes: str) -> None:
        """Write plain-text speaker notes to the slide's notes pane."""
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
